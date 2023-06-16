### info how to get annotatios text from the pdf file from this web page
### https://github.com/pymupdf/PyMuPDF/issues/819

import fitz
import os
import logging
from pprint import pprint
import pdfutils as utils
import openpyxl


### CONSTANT ###

PDF_FILE_PATH = ""
OUTPUT_FILE = "collected-comments.txt"
OUTPUT_PRESENTATION_PATH = ""
PNG_PAGE_FILE_TEMPLATE = "{path}\{prefix}-page-{page}-.png" 
IMG_LEFT = 200
IMG_TOP = 100
IMG_WIDTH = 600
IMG_HEIGTH = 600
PAR_TEXT_FONT = 10

INPUT_FILE = {
                "input1": {
                    "file_path": "input\\file1.pdf",
                    "out_path": "output\\File1-Report.ppt",
                    "csv_path": "output\\File1-Report.csv",
                    "tot_page": 160,
                    "offset_page": 0,
                    "title": "Annotazioni sezione tecnica documento file1",
                    "author": "",
                    "text_font_size": 28,
                    "page_size" : [1440, 1040]
                },
                "input2": {
                    "file_path": "input\\File2.pdf",
                    "out_path": "output\\File2-Report.ppt",
                    "csv_path": "output\\File2-Report.csv",
                    "tot_page": 46,
                    "offset_page": 0,
                    "title": "Annotazioni documento File2",
                    "author": "",
                    "text_font_size": 16,
                    "page_size" : [1440, 1040]
                },
                "input3": {
                    "file_path": "input\\File3.pdf",
                    "out_path": "output\File3-Report.ppt",
                    "csv_path": "output\File3-Report.csv",
                    "tot_page": 160,
                    "offset_page": 0,
                    "title": "Annotazioni sezione tecnica documento File3",
                    "author": "",
                    "text_font_size": 28,
                    "page_size" : [1440, 1040]
                },
                 "input4": {
                    "file_path": "input\\File4.pdf",
                    "out_path": "output\File4-Report.ppt",
                    "csv_path": "output\File4-Report.csv",
                    "tot_page": 26,
                    "offset_page": 67,
                    "title": "Annotazioni sezione tecnica documento File4",
                    "author": "",
                    "text_font_size": 28,
                    "page_size" : [1440, 1040]
                },
}

def extract_text_near_note(page,annot):
    MARGINE = 15
    try: 
        rect = annot.rect
        x0, y0, x1, y1 = rect.x0 - MARGINE, rect.y0 - MARGINE, rect.x1 + MARGINE, rect.y1 + MARGINE
        #pprint(f"x0: {x0} y0: {y0} x1: {x1} y1: {y1}")

        # New updated rec
        rect = fitz.Rect(x0, y0, x1, y1)

        page.draw_rect(rect, color=(1, 0, 0), width=2)  # Evidenzia il bordo del rettangolo con il colore rosso e larghezza 2
        text = page.get_text("text", clip=(x0, y0, x1, y1))
        return text
    except Exception as err:
        return "testo non trovato"

def measure_distance(note_rect, underlined_text_rect):
    note_x = note_rect.x0
    note_y = note_rect.y0
    
    underlined_x = underlined_text_rect.x0
    underlined_y = underlined_text_rect.y0
    
    distance = ((note_x - underlined_x)**2 + (note_y - underlined_y)**2)**0.5
    
    return distance

def extract_underlined_text_near_note(page, note):
    MARGINE = 15
    annots = page.annots()
    
    for annot in annots:
        if annot.info["content"] == '':  
            
            urect = fitz.Rect(annot.rect)  # Modifica con il rettangolo corretto del testo sottolineato
            x0, y0, x1, y1 = urect.x0 - MARGINE,urect.y0 - MARGINE, urect.x1 + MARGINE, urect.y1 + MARGINE
            urect_new = fitz.Rect(x0, y0, x1, y1)
            page.draw_rect( urect_new, color=(0, 0, 1), width=2)
            distance = measure_distance(note.rect, urect_new)

            if distance < 200:
                output = page.get_text("text", clip=urect_new)
                #print(f"Distanza: {distance} - Near text: {output}")
                return str(output)
    
    return "Not found"
    

def annotations_to_csv(file_path,ann_list):
    """annotation_to_csv

    Args:
        file_path (string): CSV file output path and filename, read fron configuration
        ann_list (list of object): [item1, item2, ...., itemn] 
                                   object itemx is a dictionary with this fields
                                   itemx = {"page": <page number>,
                                            "annotation": <string, annotation>}
    """
    try:
        logging.info("Avvio creazione file CSV")
        textfile = open(file_path, "w", newline='')  

        textfile.write("Pagina, id annotazione, Testo Vicino, Annotazione, Impatto\n")

        for ann in ann_list:  # loop through annoted pages of current PDF
            #logging.info("annuncio: %s" % ann)
            page_number = ann["page"]
            note_id = ann["note_id"]
            near_text = ann["near_text"]

            logging.info("CSV page: %d - note_id: %d - note: %s" % (page_number, note_id, str(ann["annotation"])))
       
            if type(ann["annotation"]) in (tuple, list):
                for ann_item in ann["annotation"]:
                    tmp = ann_item.replace("\n","_")
                    sub_items = tmp.partition(";")
                    textfile.write("%d, %d, %s, %s, %s, %s\n" % (page_number, note_id, near_text, sub_items[0], sub_items[1], sub_items[2])) 
            else:
                tmp = ann["annotation"].replace("\n","_")
                sub_items = tmp.partition(";")
                textfile.write("%d, %d, %s, %s, %s, %s\n" % (page_number, note_id, near_text, sub_items[0], sub_items[1], sub_items[2])) 
                     
        textfile.close() 
    except Exception as e:
        logging.error(f"Errore produzione file CSV: {e}")


def annotations_to_xlsx(file_path,ann_list):
    """annotation_to_xlsx

    Args:
        file_path (string): XLSX file output path and filename, read fron configuration
        ann_list (list of object): [item1, item2, ...., itemn] 
                                   object itemx is a dictionary with this fields
                                   itemx = {"page": <page number>,
                                            "annotation": <string, annotation>}
    """
    try:
        logging.info("Avvio creazione file XLSX")

        workbook = openpyxl.Workbook()
        worksheet = workbook.create_sheet(title='Note', index=0)

        row = 1  # Numero della riga di partenza
        col = 1
        for val in ["Pagina", "id annotazione", "Testo Vicino", "Annotazione", "Impatto"]:
            worksheet.cell(row=row, column=col, value=val)
            col = col + 1
        
        for ann in ann_list:  # loop through annoted pages of current PDF
            record = []
            page_number = ann["page"]
            note_id = ann["note_id"]
            near_text = ann["near_text"]

            record.append(page_number)
            record.append(note_id)
            record.append(near_text)

            logging.info("XLSXpage: %d - note_id: %d - note: %s" % (page_number, note_id, str(ann["annotation"])))
       
            if type(ann["annotation"]) in (tuple, list):
                for ann_item in ann["annotation"]:
                    tmp = ann_item.replace("\n","_")
                    sub_items = tmp.partition(";")
                    record.append(sub_items[0])
                    record.append(sub_items[1])
                    record.append(sub_items[2])
            else:
                tmp = ann["annotation"].replace("\n","_")
                sub_items = tmp.partition(";")
                record.append(sub_items[0])
                record.append(sub_items[1])
                record.append(sub_items[2])
            row = row + 1
            col = 1
            for val in record:
                worksheet.cell(row=row, column=col, value=val) 
                col = col + 1        
        workbook.save(file_path + ".xlsx")

    except Exception as e:
        logging.error(f"Errore produzione file XLSX: {e}")


def get_annotations(file_path,start_page, end_page, offset_page):
    """get_annotations(), retrieve all the annotation from the input pdf file

    Args:
        file_path (string): input pdf file path and filename
        start_page (int): the starting page number for the annotation extraction
        end_page (int_): the end page number for the annotation extraction_
        offset_page (int):the offset to add to calculate the document page, as written in the document

    Returns:
        list of object: the single object inside the list is a dictionary with this structure
                        { "page": <int page number>,
                          "note_id": <int number of the annotation inside the page>
                          "annotation": <string cointaining the value of the annotation>
                          "image": <string containing the path of the page's image>
                        }
    """
    
    # First check whether input file is a pdf one
    if not file_path["file_path"].endswith(".pdf"):
        logging.error("[Error] Input file is not in pdf format!!!")
        return None
    try:
        doc = fitz.open(file_path["file_path"])
    except Exception as e:
       logging.error("[Error] Error with the input file (%s)!" % str(e))
       return None 

    # Output list object
    output = []

    start_page  = start_page + offset_page - 1
    end_page = end_page + offset_page - 1
    
    #logging.info("start_page: %d - end_page: %d" % (start_page, end_page))

    for page in doc:  
        # loop through pages of current PDF
       
        #logging.info("page.number: %d", page.number)

        page_number = page.number + offset_page
      
        n_annot = 0
            
        if page_number >= start_page and page_number <= end_page:
            
            for n,annot in enumerate(page.annots()): 
                #print("-- Annotation --") 
                #pprint(dir(annot))
                #pprint(annot.info)
                
                if annot.info["content"] == '':
                    continue
       
                if annot.info["content"] == "Mio commento":
                    near_text = extract_text_near_note(page, annot)
                    #print(">>>>> Near text: %s" % near_text)

                text = annot.info["title"] + " - " + annot.info["content"]  # extract the text
                
                if len(text) != 0:
                    n_annot = n_annot +1

                    near_text = extract_underlined_text_near_note(page, annot)

                    page_annot = {
                        "page": offset_page,
                        "note_id": 1,
                        "image": "",
                        "near_text": "",
                        "annotation": [""]
                    }
                    
                    page_annot["page"] = page_number 
                    page_annot["note_id"] = n_annot
                    page_annot["near_text"] = near_text
                    #print(f"near_text: {page_annot['near_text']} - type: {type(page_annot['near_text'])}")
                
                    # Page sizes: Width & Height in pixeld
                    file_path["page_size"]= [page.rect.width, page.rect.height]

                    page_annot["annotation"] = text
                    page_annot["image"] = PNG_PAGE_FILE_TEMPLATE.format(path="output",
                                                                    prefix="tmp",
                                                                    page = page_annot["page"] )
                    #logging.info("page_number: %d - note_id: %d - text: %s" %(page_annot["page"],
                    #                                                         page_annot["note_id"],
                    #                                                          page_annot["annotation"]))
                    output.append(page_annot)                    
                else:
                    continue  
               
            utils.page_to_png(doc,page.number,page_annot["page"],"tmp","output")
    doc.save(".\\output\\output.pdf")
    doc.close()
    
    return output

####
#  list_annotations()
####
def list_annotaions(ann_list):
    output = ""
    for i,item in enumerate(ann_list):
        output = output + "({index}) {ann}\n".format(index=i, ann=item)
    return output


####
# make_presentation()
#### 
def make_presentation(file_path,ann_list,title,author,font_size, page_size):
    # Creating powerpoint presentations using the python-pptx package

    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    
    # Page ratio
    ratio = page_size[1]/page_size[0]
    logging.info("Page ratio (%d/%d): %f" % (page_size[1], page_size[0], ratio))

    # Create the Presentation object
    X = Presentation()
    X.slide_height = Pt(page_size[1])
    X.slide_width = Pt(page_size[0] * 1.5)

    # Create the presentation's front page
    layout = X.slide_layouts[0]
    first_slide = X.slides.add_slide(layout)
    first_slide.shapes.title.text = title 
    first_slide.placeholders[1].text = author
    X.save(file_path)
    
    # Set layout for the following slides
    second_layout = X.slide_layouts[5]

    # Creates following slides using annotations'list
    for item in ann_list:
        
        current_slide = X.slides.add_slide(second_layout)
        current_slide.shapes.title.text = "Note pagina #" + str(item["page"]) + " (#ann: " + str(len(item["annotation"])) + ")"
        
        # Add annotations'paragraf
        textbox = current_slide.shapes.add_textbox(Pt(page_size[0]), Pt(IMG_TOP),Pt(page_size[0]*0.5), Pt(page_size[1]*0.8)) 
        textframe = textbox.text_frame
        textframe.word_wrap = True
        textframe.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        paragraph = textframe.add_paragraph()
        paragraph.text = "%d) %s" % (item["note_id"], item["annotation"]) #list_annotaions(item["annotation"])
        paragraph.font.size = Pt(font_size)

        # Add image
        pic = current_slide.shapes.add_picture(item["image"], 
                                               Pt(10), 
                                               Pt(100), 
                                               Pt(page_size[0]*0.9),
                                               Pt(page_size[1]*0.9))

        pic.click_action.hyperlink.address = item["image"]
        X.save(file_path)


####
# Main
####

if __name__ == "__main__":
    # Setting level for logging tracing
    logging.basicConfig(level=logging.INFO)

    input_file = "input4"

    PDF_FILE_PATH = INPUT_FILE[input_file]["file_path"]
    OUTPUT_PRESENTATION_PATH = INPUT_FILE[input_file]["out_path"]
    title = INPUT_FILE[input_file]["title"]
    author = INPUT_FILE[input_file]["author"]
    font_size = INPUT_FILE[input_file]["text_font_size"]
    tot_page = INPUT_FILE[input_file]["tot_page"]
    offset_page = INPUT_FILE[input_file]["offset_page"]
   
    out = get_annotations(INPUT_FILE[input_file],1, tot_page, offset_page)
    if out != None:
        #pprint(out)
        logging.info("page size %s" % str(INPUT_FILE[input_file]["page_size"]))
        logging.info("Avvio produzione report CSV")
        annotations_to_xlsx(INPUT_FILE[input_file]["csv_path"],out)
        logging.info("Avvio produzione PPT")
        make_presentation(OUTPUT_PRESENTATION_PATH, out, title, author, font_size,INPUT_FILE[input_file]["page_size"])


